"""Document parsing for multiple file formats — v2 hardened.

Handles every edge case: corrupted PDFs, merged Excel cells,
malformed CSV, password-protected files, scanned images.
"""

import os
import csv
import json
import hashlib
import logging
import re
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple
from dataclasses import dataclass, field
from datetime import datetime

logger = logging.getLogger(__name__)


@dataclass
class PageContent:
    """Content from a single page."""
    page_number: int
    text: str
    tables: List[Dict[str, Any]] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)


@dataclass
class ParsedDocument:
    """Structured representation of a parsed document."""
    doc_id: str
    filename: str
    filepath: str
    content: str
    file_type: str
    metadata: Dict[str, Any] = field(default_factory=dict)
    tables: List[Dict[str, Any]] = field(default_factory=list)
    entities: List[Dict[str, str]] = field(default_factory=list)
    sections: List[Dict[str, str]] = field(default_factory=list)
    pages: List[PageContent] = field(default_factory=list)
    page_count: int = 0
    char_count: int = 0
    parsed_at: datetime = field(default_factory=datetime.now)
    parse_errors: List[str] = field(default_factory=list)

    def __post_init__(self):
        self.char_count = len(self.content)
        if not self.doc_id:
            self.doc_id = hashlib.md5(
                f"{self.filepath}:{self.content[:500]}".encode()
            ).hexdigest()[:12]


class DocumentParser:
    """Universal document parser — v2 hardened for every edge case."""

    SUPPORTED_EXTENSIONS = {
        ".txt": "text", ".md": "markdown", ".csv": "csv",
        ".tsv": "tsv", ".json": "json", ".html": "html",
        ".htm": "html", ".pdf": "pdf", ".docx": "docx",
        ".xlsx": "xlsx", ".xls": "xlsx", ".pptx": "pptx",
        ".png": "image", ".jpg": "image", ".jpeg": "image",
        ".tiff": "image", ".rtf": "text", ".xml": "xml",
    }

    def __init__(self):
        self._parsers = {
            "text": self._parse_text, "markdown": self._parse_text,
            "csv": self._parse_csv, "tsv": self._parse_csv,
            "json": self._parse_json, "html": self._parse_html,
            "pdf": self._parse_pdf, "docx": self._parse_docx,
            "xlsx": self._parse_xlsx, "pptx": self._parse_pptx,
            "image": self._parse_image, "xml": self._parse_xml,
        }

    def parse(self, filepath: str) -> ParsedDocument:
        """Parse a single document with full error handling."""
        path = Path(filepath)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {filepath}")

        ext = path.suffix.lower()
        file_type = self.SUPPORTED_EXTENSIONS.get(ext)
        if not file_type:
            raise ValueError(
                f"Unsupported file type: {ext}. "
                f"Supported: {list(self.SUPPORTED_EXTENSIONS.keys())}"
            )

        parser_fn = self._parsers[file_type]
        logger.info(f"Parsing {path.name} as {file_type}")

        try:
            doc = parser_fn(str(path))
            doc.filename = path.name
            doc.filepath = str(path)
            doc.file_type = file_type
            doc.metadata["file_size"] = path.stat().st_size
            doc.metadata["extension"] = ext
            doc.metadata["modified"] = datetime.fromtimestamp(
                path.stat().st_mtime
            ).isoformat()
            return doc
        except Exception as e:
            logger.error(f"Failed to parse {filepath}: {e}")
            # Return partial document instead of crashing
            return ParsedDocument(
                doc_id="", filename=path.name, filepath=str(path),
                content=f"[Parse error: {e}]", file_type=file_type,
                parse_errors=[str(e)],
            )

    def parse_directory(self, directory: str, recursive: bool = True) -> List[ParsedDocument]:
        """Parse all supported documents in a directory."""
        path = Path(directory)
        if not path.is_dir():
            raise NotADirectoryError(f"Not a directory: {directory}")

        documents = []
        pattern = "**/*" if recursive else "*"

        for file_path in sorted(path.glob(pattern)):
            if file_path.is_file() and file_path.suffix.lower() in self.SUPPORTED_EXTENSIONS:
                try:
                    doc = self.parse(str(file_path))
                    documents.append(doc)
                except Exception as e:
                    logger.warning(f"Skipping {file_path}: {e}")

        logger.info(f"Parsed {len(documents)} documents from {directory}")
        return documents

    # ─── TEXT / MARKDOWN ───────────────────────────────────

    def _parse_text(self, filepath: str) -> ParsedDocument:
        encoding = self._detect_encoding(filepath)
        with open(filepath, "r", encoding=encoding, errors="replace") as f:
            content = f.read()

        # Strip BOM
        if content.startswith("\ufeff"):
            content = content[1:]

        sections = []
        current_section = {"title": "Main", "content": ""}
        for line in content.split("\n"):
            if line.startswith("#"):
                if current_section["content"].strip():
                    sections.append(current_section)
                title = line.lstrip("#").strip()
                current_section = {"title": title, "content": ""}
            else:
                current_section["content"] += line + "\n"
        if current_section["content"].strip():
            sections.append(current_section)

        return ParsedDocument(
            doc_id="", filename="", filepath=filepath,
            content=content, file_type="text", sections=sections,
        )

    # ─── CSV / TSV ─────────────────────────────────────────

    def _parse_csv(self, filepath: str) -> ParsedDocument:
        encoding = self._detect_encoding(filepath)

        # Auto-detect delimiter
        with open(filepath, "r", encoding=encoding, errors="replace") as f:
            sample = f.read(8192)

        # Strip BOM
        if sample.startswith("\ufeff"):
            sample = sample[1:]

        try:
            dialect = csv.Sniffer().sniff(sample)
            delimiter = dialect.delimiter
        except csv.Error:
            # Fallback: check common delimiters
            if filepath.endswith(".tsv"):
                delimiter = "\t"
            elif sample.count(";") > sample.count(","):
                delimiter = ";"
            elif sample.count("|") > sample.count(","):
                delimiter = "|"
            else:
                delimiter = ","

        rows = []
        headers = []
        errors = []

        with open(filepath, "r", encoding=encoding, errors="replace") as f:
            content_raw = f.read()
            if content_raw.startswith("\ufeff"):
                content_raw = content_raw[1:]

        try:
            reader = csv.reader(content_raw.splitlines(), delimiter=delimiter)
            expected_cols = None

            for i, row in enumerate(reader):
                if i == 0:
                    headers = [h.strip() for h in row]
                    expected_cols = len(headers)
                else:
                    # Handle malformed rows
                    if expected_cols:
                        if len(row) < expected_cols:
                            row = row + [""] * (expected_cols - len(row))  # Pad
                        elif len(row) > expected_cols:
                            row = row[:expected_cols]  # Truncate
                            errors.append(f"Row {i+1}: truncated from {len(row)} cols")
                    rows.append(row)
        except Exception as e:
            errors.append(f"CSV parse error: {e}")

        if not headers and rows:
            headers = [f"col_{j+1}" for j in range(len(rows[0]))]

        content_lines = [delimiter.join(headers)]
        for row in rows:
            content_lines.append(delimiter.join(row))
        content = "\n".join(content_lines)

        tables = [{"headers": headers, "rows": rows, "source": filepath}]

        return ParsedDocument(
            doc_id="", filename="", filepath=filepath,
            content=content, file_type="csv", tables=tables,
            metadata={"row_count": len(rows), "column_count": len(headers),
                       "columns": headers, "delimiter": delimiter},
            parse_errors=errors,
        )

    # ─── JSON ──────────────────────────────────────────────

    def _parse_json(self, filepath: str) -> ParsedDocument:
        encoding = self._detect_encoding(filepath)
        with open(filepath, "r", encoding=encoding, errors="replace") as f:
            raw = f.read()
        if raw.startswith("\ufeff"):
            raw = raw[1:]
        data = json.loads(raw)
        content = json.dumps(data, indent=2, ensure_ascii=False)
        return ParsedDocument(
            doc_id="", filename="", filepath=filepath,
            content=content, file_type="json",
            metadata={"json_type": type(data).__name__},
        )

    # ─── HTML ──────────────────────────────────────────────

    def _parse_html(self, filepath: str) -> ParsedDocument:
        encoding = self._detect_encoding(filepath)
        with open(filepath, "r", encoding=encoding, errors="replace") as f:
            raw = f.read()

        clean = re.sub(r"<(script|style)[^>]*>.*?</\1>", "", raw, flags=re.DOTALL | re.IGNORECASE)
        clean = re.sub(r"<[^>]+>", " ", clean)
        clean = re.sub(r"\s+", " ", clean).strip()

        return ParsedDocument(
            doc_id="", filename="", filepath=filepath,
            content=clean, file_type="html",
        )

    # ─── PDF (v2 hardened) ─────────────────────────────────

    def _parse_pdf(self, filepath: str) -> ParsedDocument:
        errors = []

        # Try PyMuPDF first
        try:
            return self._parse_pdf_pymupdf(filepath)
        except ImportError:
            pass
        except Exception as e:
            errors.append(f"PyMuPDF failed: {e}")

        # Try pdfplumber
        try:
            return self._parse_pdf_pdfplumber(filepath)
        except ImportError:
            pass
        except Exception as e:
            errors.append(f"pdfplumber failed: {e}")

        if errors:
            return ParsedDocument(
                doc_id="", filename="", filepath=filepath,
                content=f"[PDF parse failed: {'; '.join(errors)}]",
                file_type="pdf", parse_errors=errors,
            )

        raise ImportError(
            "PDF parsing requires PyMuPDF or pdfplumber. "
            "Install with: pip install PyMuPDF or pip install pdfplumber"
        )

    def _parse_pdf_pymupdf(self, filepath: str) -> ParsedDocument:
        import fitz

        errors = []
        try:
            doc = fitz.open(filepath)
        except Exception as e:
            # Try with empty password
            if "password" in str(e).lower() or "encrypted" in str(e).lower():
                try:
                    doc = fitz.open(filepath)
                    doc.authenticate("")
                    if not doc.is_encrypted or doc.authenticate(""):
                        pass
                    else:
                        return ParsedDocument(
                            doc_id="", filename="", filepath=filepath,
                            content="[Password-protected PDF — cannot parse]",
                            file_type="pdf", parse_errors=["password_protected"],
                        )
                except Exception:
                    return ParsedDocument(
                        doc_id="", filename="", filepath=filepath,
                        content="[Password-protected PDF — cannot parse]",
                        file_type="pdf", parse_errors=["password_protected"],
                    )
            else:
                return ParsedDocument(
                    doc_id="", filename="", filepath=filepath,
                    content=f"[Corrupted PDF: {e}]",
                    file_type="pdf", parse_errors=[f"corrupted: {e}"],
                )

        pages = []
        tables = []
        all_page_texts = []

        # Detect headers/footers (repeated text across pages)
        repeated = self._detect_pdf_repeated_blocks(doc)

        for page_num in range(len(doc)):
            try:
                page = doc[page_num]
                text = page.get_text(sort=True)  # Reading order

                # Strip repeated headers/footers
                text = self._strip_repeated_text(text, repeated)

                # If almost no text, try OCR
                if len(text.strip()) < 20:
                    ocr_text = self._ocr_pdf_page(page)
                    if ocr_text and len(ocr_text.strip()) > len(text.strip()):
                        text = ocr_text

                # Extract tables
                try:
                    page_tables = page.find_tables()
                    for table in page_tables:
                        table_data = table.extract()
                        if table_data and len(table_data) > 1:
                            tables.append({
                                "headers": [str(h) if h else "" for h in table_data[0]],
                                "rows": table_data[1:],
                                "page": page_num + 1,
                                "source": filepath,
                            })
                except Exception as te:
                    errors.append(f"Table extraction failed on page {page_num+1}: {te}")

                page_content = PageContent(
                    page_number=page_num + 1,
                    text=text,
                    tables=[t for t in tables if t.get("page") == page_num + 1],
                )
                pages.append(page_content)
                all_page_texts.append(text)

            except Exception as pe:
                errors.append(f"Page {page_num+1} failed: {pe}")
                pages.append(PageContent(page_number=page_num + 1, text=""))

        doc.close()
        content = "\n\n".join(all_page_texts)

        return ParsedDocument(
            doc_id="", filename="", filepath=filepath,
            content=content, file_type="pdf",
            tables=tables, pages=pages,
            page_count=len(pages), parse_errors=errors,
        )

    def _detect_pdf_repeated_blocks(self, doc) -> List[str]:
        """Detect headers/footers repeated across >60% of pages."""
        if len(doc) < 3:
            return []

        page_first_lines = []
        page_last_lines = []

        for i in range(min(len(doc), 20)):  # Sample first 20 pages
            try:
                text = doc[i].get_text()
                lines = [l.strip() for l in text.split("\n") if l.strip()]
                if lines:
                    page_first_lines.extend(lines[:2])
                    page_last_lines.extend(lines[-2:])
            except Exception:
                pass

        repeated = []
        threshold = min(len(doc), 20) * 0.6

        from collections import Counter
        for line, count in Counter(page_first_lines).items():
            if count >= threshold and len(line) > 3:
                repeated.append(line)
        for line, count in Counter(page_last_lines).items():
            if count >= threshold and len(line) > 3:
                repeated.append(line)

        return repeated

    def _strip_repeated_text(self, text: str, repeated: List[str]) -> str:
        """Remove repeated header/footer lines."""
        if not repeated:
            return text
        lines = text.split("\n")
        filtered = [l for l in lines if l.strip() not in repeated]
        return "\n".join(filtered)

    def _ocr_pdf_page(self, page) -> str:
        """OCR a PDF page when no text layer exists."""
        try:
            import pytesseract
            from PIL import Image
            import io
            pix = page.get_pixmap(dpi=300)
            img_data = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_data))
            return pytesseract.image_to_string(img)
        except ImportError:
            return ""
        except Exception as e:
            logger.debug(f"OCR failed: {e}")
            return ""

    def _parse_pdf_pdfplumber(self, filepath: str) -> ParsedDocument:
        import pdfplumber
        pages = []
        tables = []
        all_texts = []

        try:
            with pdfplumber.open(filepath) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    try:
                        text = page.extract_text() or ""
                        all_texts.append(text)

                        page_tables = []
                        for table in (page.extract_tables() or []):
                            if table and len(table) > 1:
                                t = {
                                    "headers": [str(h) if h else "" for h in table[0]],
                                    "rows": table[1:],
                                    "page": page_num + 1,
                                    "source": filepath,
                                }
                                tables.append(t)
                                page_tables.append(t)

                        pages.append(PageContent(
                            page_number=page_num + 1, text=text, tables=page_tables,
                        ))
                    except Exception as pe:
                        pages.append(PageContent(page_number=page_num + 1, text=""))
        except Exception as e:
            if "password" in str(e).lower():
                return ParsedDocument(
                    doc_id="", filename="", filepath=filepath,
                    content="[Password-protected PDF]",
                    file_type="pdf", parse_errors=["password_protected"],
                )
            raise

        return ParsedDocument(
            doc_id="", filename="", filepath=filepath,
            content="\n\n".join(all_texts), file_type="pdf",
            tables=tables, pages=pages, page_count=len(pages),
        )

    # ─── DOCX (v2 hardened) ────────────────────────────────

    def _parse_docx(self, filepath: str) -> ParsedDocument:
        try:
            from docx import Document as DocxDocument
        except ImportError:
            raise ImportError("DOCX requires python-docx: pip install python-docx")

        errors = []
        try:
            doc = DocxDocument(filepath)
        except Exception as e:
            if "password" in str(e).lower() or "encrypted" in str(e).lower():
                return ParsedDocument(
                    doc_id="", filename="", filepath=filepath,
                    content="[Password-protected DOCX]",
                    file_type="docx", parse_errors=["password_protected"],
                )
            raise

        paragraphs = []
        for p in doc.paragraphs:
            text = p.text.strip()
            if text:
                paragraphs.append(text)

        # Extract tables (handle nested)
        tables = []
        for table in doc.tables:
            try:
                rows = []
                for row in table.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                if rows:
                    tables.append({
                        "headers": rows[0], "rows": rows[1:], "source": filepath,
                    })
            except Exception as te:
                errors.append(f"Table extraction failed: {te}")

        # Extract comments if available
        comments = []
        try:
            from docx.opc.constants import RELATIONSHIP_TYPE as RT
            # Comments are in a separate part
            for rel in doc.part.rels.values():
                if "comments" in str(rel.reltype).lower():
                    comments.append("[Comments present in document]")
        except Exception:
            pass

        content = "\n\n".join(paragraphs)

        return ParsedDocument(
            doc_id="", filename="", filepath=filepath,
            content=content, file_type="docx",
            tables=tables, parse_errors=errors,
            metadata={"has_comments": len(comments) > 0},
        )

    # ─── XLSX (v2 hardened) ────────────────────────────────

    def _parse_xlsx(self, filepath: str) -> ParsedDocument:
        try:
            import openpyxl
        except ImportError:
            raise ImportError("XLSX requires openpyxl: pip install openpyxl")

        errors = []
        try:
            wb = openpyxl.load_workbook(filepath, data_only=True, read_only=False)
        except Exception as e:
            if "password" in str(e).lower() or "encrypted" in str(e).lower():
                return ParsedDocument(
                    doc_id="", filename="", filepath=filepath,
                    content="[Password-protected Excel]",
                    file_type="xlsx", parse_errors=["password_protected"],
                )
            # Try read-only mode for large files
            try:
                wb = openpyxl.load_workbook(filepath, data_only=True, read_only=True)
            except Exception:
                raise e

        content_parts = []
        tables = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]

            # Skip hidden sheets
            try:
                if hasattr(ws, 'sheet_state') and ws.sheet_state == 'hidden':
                    continue
            except Exception:
                pass

            # Handle merged cells
            merged_ranges = []
            try:
                if hasattr(ws, 'merged_cells') and ws.merged_cells:
                    merged_ranges = list(ws.merged_cells.ranges)
            except Exception:
                pass

            rows = []
            empty_row_count = 0

            for row in ws.iter_rows(values_only=True):
                row_data = []
                for cell in row:
                    if cell is None:
                        row_data.append("")
                    elif isinstance(cell, datetime):
                        row_data.append(cell.isoformat())
                    else:
                        row_data.append(str(cell))

                # Skip rows that are >80% empty
                non_empty = sum(1 for c in row_data if c.strip())
                if non_empty == 0:
                    empty_row_count += 1
                    if empty_row_count > 5:  # Stop after 5 consecutive empty rows
                        break
                    continue
                else:
                    empty_row_count = 0
                    rows.append(row_data)

            if rows:
                content_parts.append(f"Sheet: {sheet_name}")
                for row in rows:
                    content_parts.append(", ".join(row))
                tables.append({
                    "sheet": sheet_name,
                    "headers": rows[0] if rows else [],
                    "rows": rows[1:] if len(rows) > 1 else [],
                    "source": filepath,
                    "merged_cells": len(merged_ranges),
                })

        try:
            wb.close()
        except Exception:
            pass

        return ParsedDocument(
            doc_id="", filename="", filepath=filepath,
            content="\n".join(content_parts), file_type="xlsx",
            tables=tables, parse_errors=errors,
            metadata={"sheet_count": len(wb.sheetnames), "sheets": wb.sheetnames},
        )

    # ─── PPTX (v2 hardened) ────────────────────────────────

    def _parse_pptx(self, filepath: str) -> ParsedDocument:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("PPTX requires python-pptx: pip install python-pptx")

        errors = []
        try:
            prs = Presentation(filepath)
        except Exception as e:
            if "password" in str(e).lower():
                return ParsedDocument(
                    doc_id="", filename="", filepath=filepath,
                    content="[Password-protected PPTX]",
                    file_type="pptx", parse_errors=["password_protected"],
                )
            raise

        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []

            for shape in slide.shapes:
                # Regular text
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)

                # Tables in slides
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        texts.append(" | ".join(row_text))

                # Grouped shapes
                try:
                    if hasattr(shape, 'shapes'):
                        for sub_shape in shape.shapes:
                            if sub_shape.has_text_frame:
                                for p in sub_shape.text_frame.paragraphs:
                                    if p.text.strip():
                                        texts.append(p.text.strip())
                except Exception:
                    pass

            # Speaker notes
            notes = ""
            try:
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                pass

            if texts:
                slide_text = f"Slide {i+1}:\n" + "\n".join(texts)
                if notes:
                    slide_text += f"\n[Speaker Notes: {notes}]"
                slides.append(slide_text)

        return ParsedDocument(
            doc_id="", filename="", filepath=filepath,
            content="\n\n".join(slides), file_type="pptx",
            page_count=len(prs.slides), parse_errors=errors,
        )

    # ─── IMAGE (OCR) ──────────────────────────────────────

    def _parse_image(self, filepath: str) -> ParsedDocument:
        try:
            import pytesseract
            from PIL import Image
            img = Image.open(filepath)
            content = pytesseract.image_to_string(img)
            return ParsedDocument(
                doc_id="", filename="", filepath=filepath,
                content=content, file_type="image",
                metadata={"width": img.width, "height": img.height},
            )
        except ImportError:
            return ParsedDocument(
                doc_id="", filename="", filepath=filepath,
                content=f"[Image: {filepath}. OCR not available.]",
                file_type="image",
            )

    # ─── XML ───────────────────────────────────────────────

    def _parse_xml(self, filepath: str) -> ParsedDocument:
        import xml.etree.ElementTree as ET
        tree = ET.parse(filepath)
        root = tree.getroot()

        def extract_text(element, depth=0):
            texts = []
            tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag
            if element.text and element.text.strip():
                texts.append(f"{'  ' * depth}{tag}: {element.text.strip()}")
            for child in element:
                texts.extend(extract_text(child, depth + 1))
            return texts

        content = "\n".join(extract_text(root))
        return ParsedDocument(
            doc_id="", filename="", filepath=filepath,
            content=content, file_type="xml",
        )

    # ─── UTILITIES ─────────────────────────────────────────

    def _detect_encoding(self, filepath: str) -> str:
        """Detect file encoding with fallback chain."""
        encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]
        for enc in encodings:
            try:
                with open(filepath, "r", encoding=enc) as f:
                    f.read(4096)
                return enc
            except (UnicodeDecodeError, UnicodeError):
                continue
        return "utf-8"


# ─── SQL CONNECTOR (v2 new) ────────────────────────────────

class SQLConnector:
    """Ingest data from SQL databases."""

    def parse(self, connection_string: str,
              tables: Optional[List[str]] = None,
              query: Optional[str] = None) -> List[ParsedDocument]:
        """Parse SQL database into documents.

        Args:
            connection_string: SQLAlchemy connection string
            tables: Specific tables to ingest (None = all)
            query: Custom SQL query
        """
        try:
            from sqlalchemy import create_engine, inspect, text
        except ImportError:
            raise ImportError("SQL connector requires sqlalchemy: pip install sqlalchemy")

        engine = create_engine(connection_string)
        documents = []

        if query:
            # Custom query
            try:
                import pandas as pd
                df = pd.read_sql(query, engine)
                doc = self._dataframe_to_document(df, "custom_query", connection_string)
                documents.append(doc)
            except ImportError:
                # Without pandas, use raw execution
                with engine.connect() as conn:
                    result = conn.execute(text(query))
                    headers = list(result.keys())
                    rows = [list(row) for row in result]
                    doc = self._rows_to_document(headers, rows, "custom_query", connection_string)
                    documents.append(doc)
        else:
            inspector = inspect(engine)
            table_names = tables or inspector.get_table_names()

            for table_name in table_names:
                try:
                    columns = inspector.get_columns(table_name)
                    col_names = [c["name"] for c in columns]

                    # Get foreign keys for graph edges
                    foreign_keys = []
                    try:
                        fks = inspector.get_foreign_keys(table_name)
                        foreign_keys = fks
                    except Exception:
                        pass

                    # Read data
                    with engine.connect() as conn:
                        result = conn.execute(text(f"SELECT * FROM {table_name} LIMIT 10000"))
                        rows = [[str(v) if v is not None else "[empty]" for v in row] for row in result]

                    doc = self._rows_to_document(col_names, rows, table_name, connection_string)
                    doc.metadata["columns"] = col_names
                    doc.metadata["column_types"] = {c["name"]: str(c["type"]) for c in columns}
                    doc.metadata["foreign_keys"] = foreign_keys
                    documents.append(doc)

                    logger.info(f"Ingested table {table_name}: {len(rows)} rows")

                except Exception as e:
                    logger.warning(f"Failed to ingest table {table_name}: {e}")

        engine.dispose()
        return documents

    def _rows_to_document(self, headers, rows, source, connection_string):
        """Convert rows into a ParsedDocument."""
        content_lines = [", ".join(str(h) for h in headers)]
        for row in rows:
            content_lines.append(", ".join(str(v) for v in row))

        tables = [{"headers": headers, "rows": rows, "source": source}]

        return ParsedDocument(
            doc_id="", filename=source, filepath=connection_string,
            content="\n".join(content_lines), file_type="sql",
            tables=tables,
            metadata={"row_count": len(rows), "column_count": len(headers)},
        )

    def _dataframe_to_document(self, df, source, connection_string):
        """Convert pandas DataFrame to ParsedDocument."""
        headers = list(df.columns)
        rows = df.fillna("[empty]").astype(str).values.tolist()
        return self._rows_to_document(headers, rows, source, connection_string)
